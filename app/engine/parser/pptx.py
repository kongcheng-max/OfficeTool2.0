"""PPTX 解析器 — 基于 python-pptx（线程池执行，避免阻塞 event loop）"""

from typing import List

from pptx import Presentation

from engine.parser.base import BaseParser, Chunk


class PPTXParser(BaseParser):
    name = "pptx"
    supported_extensions = [".pptx"]
    supported_mime_types = [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]

    async def parse(self, file_path: str, original_filename: str) -> List[Chunk]:
        """PPTX 解析 — python-pptx + lxml C 扩展同步调用，用线程池隔离"""
        return await self._run_sync_in_thread(self._parse_sync, file_path, original_filename)

    def _parse_sync(self, file_path: str, original_filename: str) -> List[Chunk]:
        chunks: List[Chunk] = []
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            parts = []

            # 标题
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
                parts.append(f"# {title}")

            # 正文（文本框）
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    text = shape.text.strip()
                    if text:
                        parts.append(text)

                # 表格
                if shape.has_table:
                    table = shape.table
                    table_md = self._table_to_markdown(table)
                    if table_md.strip():
                        parts.append(table_md)

            # 备注
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"\n> 备注: {notes}")

            if parts:
                content = "\n\n".join(parts)
                chunks.append(Chunk(
                    content=content,
                    metadata={
                        "slide": slide_num,
                        "source": original_filename,
                        "parser_name": self.name,
                        "chunk_index": len(chunks),
                    },
                    chunk_type="text",
                ))

        return chunks

    @staticmethod
    def _table_to_markdown(table) -> str:
        rows = table.rows
        if not rows:
            return ""
        lines = []
        header_cells = [cell.text.strip() for cell in rows[0].cells]
        lines.append("| " + " | ".join(header_cells) + " |")
        lines.append("|" + "|".join("---" for _ in header_cells) + "|")
        for row in rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
        return "\n".join(lines)
